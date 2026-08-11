"""
File Processing Service
Handles document reading, text extraction, and chunking
Supports: PDF, DOCX, PPTX, TXT
"""

import os
import zipfile
import io
from pathlib import Path
from typing import List, Tuple, Optional, Dict
import uuid
from datetime import datetime
import logging
from PIL import Image
import google.generativeai as genai

from pypdf import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation
import re

from app.utils.config import settings, UPLOAD_DIR
from app.utils.logger import get_logger
from app.models.schemas import TextChunk, StoredFileMetadata

logger = get_logger(__name__)


class FileValidator:
    
    ALLOWED_EXTENSIONS = {"pdf", "docx", "pptx", "txt"}
    MAX_FILE_SIZE = settings.max_file_size
    
    @staticmethod
    def validate_file(filename: str, file_size: int) -> Tuple[bool, Optional[str]]:
        file_ext = filename.split('.')[-1].lower()
        if file_ext not in FileValidator.ALLOWED_EXTENSIONS:
            return False, f"Invalid file type: {file_ext}. Allowed: {', '.join(FileValidator.ALLOWED_EXTENSIONS)}"
        
        if file_size > FileValidator.MAX_FILE_SIZE:
            max_mb = FileValidator.MAX_FILE_SIZE / (1024 * 1024)
            file_mb = file_size / (1024 * 1024)
            return False, f"File too large: {file_mb:.1f}MB (max: {max_mb:.1f}MB)"
        
        logger.info(f"✅ File validation passed: {filename} ({file_size / 1024:.1f}KB)")
        return True, None


class TextExtractor:
    
    @staticmethod
    def extract_from_pdf(file_path: str) -> Tuple[str, int, List[str]]:
        logger.info(f"Extracting text from PDF: {file_path}")
        
        try:
            pdf_reader = PdfReader(file_path)
            num_pages = len(pdf_reader.pages)
            page_texts = []
            full_text = ""
            
            for page_num, page in enumerate(pdf_reader.pages):
                page_text = page.extract_text() or ""
                page_texts.append(page_text)
                full_text += f"\n\n--- Page {page_num + 1} ---\n{page_text}"
            
            logger.info(f"✅ Extracted {num_pages} pages from PDF")
            return full_text, num_pages, page_texts
            
        except Exception as e:
            logger.error(f"❌ Error extracting PDF: {str(e)}")
            raise Exception(f"Failed to extract PDF: {str(e)}")

    @staticmethod
    def extract_images_from_docx(file_path: str) -> str:
        extracted_text = ""
        try:
            with zipfile.ZipFile(file_path) as z:
                media_files = sorted([n for n in z.namelist() if n.startswith("word/media/") and n.lower().endswith((".png", ".jpg", ".jpeg", ".bmp", ".webp"))])
                if not media_files:
                    return ""
                
                images = []
                for name in media_files:
                    try:
                        img_bytes = z.read(name)
                        img = Image.open(io.BytesIO(img_bytes)).convert("RGB")
                        images.append(img)
                    except Exception as img_err:
                        logger.warning(f"Could not load image {name}: {img_err}")
                
                if not images:
                    return ""

                genai.configure(api_key=settings.google_api_key)
                models_to_try = ["gemini-3.5-flash-lite", settings.gemini_model, "gemini-flash-latest", "gemini-3.6-flash", "gemini-2.5-flash"]
                
                prompt = ["Extract all readable text, code, questions, tasks, and contents from these document images in order, exactly as written:"] + images
                
                for m_name in models_to_try:
                    if not m_name:
                        continue
                    try:
                        model = genai.GenerativeModel(m_name)
                        resp = model.generate_content(prompt)
                        if resp and resp.text:
                            extracted_text = resp.text.strip()
                            logger.info(f"✅ Image OCR succeeded with model {m_name} ({len(extracted_text)} chars)")
                            break
                    except Exception as m_err:
                        logger.warning(f"Model {m_name} failed for document image batch: {m_err}")
        except Exception as e:
            logger.warning(f"Failed docx image fallback: {e}")
        
        return extracted_text.strip()

    @staticmethod
    def extract_from_docx(file_path: str) -> Tuple[str, int, List[str]]:
        logger.info(f"Extracting text from DOCX: {file_path}")
        
        try:
            doc = DocxDocument(file_path)
            full_text = ""
            paragraph_texts = []
            
            for para_num, para in enumerate(doc.paragraphs):
                if para.text.strip():
                    paragraph_texts.append(para.text)
                    full_text += f"\n{para.text}"
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells if cell.text.strip()])
                    if row_text.strip():
                        paragraph_texts.append(row_text)
                        full_text += f"\n{row_text}"
            
            if len(full_text.strip()) < 20:
                logger.info(f"Plain text in DOCX is short or empty ({len(full_text.strip())} chars), running image OCR fallback...")
                image_text = TextExtractor.extract_images_from_docx(file_path)
                if image_text:
                    full_text += f"\n{image_text}"
                    paragraph_texts.extend([line for line in image_text.split('\n') if line.strip()])
            
            logger.info(f"✅ Extracted {len(paragraph_texts)} text blocks ({len(full_text)} chars) from DOCX")
            return full_text, len(paragraph_texts), paragraph_texts
            
        except Exception as e:
            logger.error(f"❌ Error extracting DOCX: {str(e)}")
            raise Exception(f"Failed to extract DOCX: {str(e)}")
    
    @staticmethod
    def extract_from_pptx(file_path: str) -> Tuple[str, int, List[str]]:
        """
        Extract text from PPTX (PowerPoint) file
        
        Args:
            file_path: Path to PPTX file
            
        Returns:
            Tuple[str, int, List[str]]: (full_text, num_slides, slide_texts)
        """
        logger.info(f"Extracting text from PPTX: {file_path}")
        
        try:
            presentation = Presentation(file_path)
            full_text = ""
            slide_texts = []
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_text = ""
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += f"\n{shape.text}"
                
                slide_texts.append(slide_text)
                full_text += f"\n\n--- Slide {slide_num + 1} ---\n{slide_text}"
            
            logger.info(f"✅ Extracted {len(presentation.slides)} slides from PPTX")
            return full_text, len(presentation.slides), slide_texts
            
        except Exception as e:
            logger.error(f"❌ Error extracting PPTX: {str(e)}")
            raise Exception(f"Failed to extract PPTX: {str(e)}")
    
    @staticmethod
    def extract_from_txt(file_path: str) -> Tuple[str, int, List[str]]:
        """
        Extract text from TXT file
        
        Args:
            file_path: Path to TXT file
            
        Returns:
            Tuple[str, int, List[str]]: (full_text, num_lines, line_texts)
        """
        logger.info(f"Extracting text from TXT: {file_path}")
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                full_text = f.read()
            
            line_texts = [line.strip() for line in full_text.split('\n') if line.strip()]
            
            logger.info(f"✅ Extracted {len(line_texts)} lines from TXT")
            return full_text, len(line_texts), line_texts
            
        except Exception as e:
            logger.error(f"❌ Error extracting TXT: {str(e)}")
            raise Exception(f"Failed to extract TXT: {str(e)}")
    
    @staticmethod
    def extract_text(file_path: str, file_type: str) -> Tuple[str, int, List[str]]:
        """
        Route to appropriate extraction method based on file type
        
        Args:
            file_path: Path to file
            file_type: File type (pdf, docx, pptx, txt)
            
        Returns:
            Tuple[str, int, List[str]]: (full_text, count, details)
        """
        file_type = file_type.lower()
        
        if file_type == "pdf":
            return TextExtractor.extract_from_pdf(file_path)
        elif file_type == "docx":
            return TextExtractor.extract_from_docx(file_path)
        elif file_type == "pptx":
            return TextExtractor.extract_from_pptx(file_path)
        elif file_type == "txt":
            return TextExtractor.extract_from_txt(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")


# ============ TEXT CHUNKING ============

class TextChunker:
    """Splits text into overlapping chunks"""
    
    @staticmethod
    def clean_text(text: str) -> str:
        if not text:
            return ""
        if len(text) > 500000:
            return text.replace('\r', ' ').replace('\t', ' ').strip()
        return re.sub(r'[ \t\r\f\v]+', ' ', text).strip()
    
    @staticmethod
    def chunk_text(
        text: str,
        chunk_size: int = None,
        overlap: int = None
    ) -> List[TextChunk]:
        total_len = len(text)
        if chunk_size is None:
            if total_len < 100000:
                chunk_size = 2000
                overlap = 200
            elif total_len < 1000000:
                chunk_size = 5000
                overlap = 300
            else:
                chunk_size = 10000
                overlap = 500
        elif overlap is None:
            overlap = settings.chunk_overlap
        
        logger.info(f"Chunking text: len={total_len}, chunk_size={chunk_size}, overlap={overlap}")
        
        text = TextChunker.clean_text(text)
        
        if len(text) < chunk_size:
            chunk = TextChunk(
                chunk_id=str(uuid.uuid4()),
                file_id="",
                chunk_index=0,
                text=text,
                page_number=None,
                section_heading=None
            )
            logger.info("✅ Created 1 chunk (text smaller than chunk_size)")
            return [chunk]
        
        chunks = []
        start_idx = 0
        chunk_index = 0
        
        while start_idx < len(text):
            end_idx = min(start_idx + chunk_size, len(text))
            chunk_text = text[start_idx:end_idx].strip()
            
            if chunk_text:
                chunk = TextChunk(
                    chunk_id=str(uuid.uuid4()),
                    file_id="",
                    chunk_index=chunk_index,
                    text=chunk_text,
                    page_number=None,
                    section_heading=None
                )
                chunks.append(chunk)
                chunk_index += 1
            
            start_idx = end_idx - overlap
            if start_idx >= end_idx or end_idx >= len(text):
                break
        
        logger.info(f"✅ Created {len(chunks)} chunks from text")
        return chunks
    
    @staticmethod
    def detect_page_numbers(chunks: List[TextChunk], page_markers: List[str]) -> List[TextChunk]:
        current_page = None
        for chunk in chunks:
            match = re.search(r'Page (\d+)', chunk.text)
            if match:
                try:
                    current_page = int(match.group(1))
                except ValueError:
                    pass
            if current_page:
                chunk.page_number = current_page
        return chunks


# ============ FILE STORAGE ============

class FileStorage:
    """Manages file storage on disk"""
    
    @staticmethod
    def save_uploaded_file(
        file_content: bytes,
        filename: str,
        collection_id: str
    ) -> Tuple[str, str]:
        """
        Save uploaded file to disk
        
        Args:
            file_content: Binary file content
            filename: Original filename
            collection_id: User/collection ID
            
        Returns:
            Tuple[str, str]: (file_path, file_id)
        """
        # Create directory for collection if it doesn't exist
        collection_dir = UPLOAD_DIR / collection_id
        collection_dir.mkdir(parents=True, exist_ok=True)
        
        # Generate unique file ID
        file_id = str(uuid.uuid4())
        file_ext = filename.split('.')[-1].lower()
        
        # Create unique filename
        unique_filename = f"{file_id}.{file_ext}"
        file_path = collection_dir / unique_filename
        
        # Save file
        with open(file_path, 'wb') as f:
            f.write(file_content)
        
        logger.info(f"✅ File saved: {file_path} (ID: {file_id})")
        return str(file_path), file_id
    
    @staticmethod
    def delete_file(file_path: str) -> bool:
        """
        Delete file from disk
        
        Args:
            file_path: Path to file
            
        Returns:
            bool: Success status
        """
        try:
            path = Path(file_path)
            if path.exists():
                path.unlink()
                logger.info(f"✅ File deleted: {file_path}")
                return True
            else:
                logger.warning(f"⚠️ File not found for deletion: {file_path}")
                return False
        except Exception as e:
            logger.error(f"❌ Error deleting file: {str(e)}")
            return False


# ============ MAIN FILE PROCESSOR ============

class FileProcessor:
    """Main orchestrator for file processing"""
    
    def __init__(self):
        self.validator = FileValidator()
        self.extractor = TextExtractor()
        self.chunker = TextChunker()
        self.storage = FileStorage()
    
    def process_file(
        self,
        file_content: bytes,
        filename: str,
        file_type: str,
        collection_id: str = "default"
    ) -> Tuple[List[TextChunk], StoredFileMetadata]:
        """
        Complete file processing pipeline
        
        Pipeline:
        1. Validate file
        2. Save to disk
        3. Extract text
        4. Chunk text
        5. Return chunks + metadata
        
        Args:
            file_content: Binary file content
            filename: Original filename
            file_type: File type (pdf, docx, pptx, txt)
            collection_id: User/collection ID
            
        Returns:
            Tuple[List[TextChunk], StoredFileMetadata]: (chunks, metadata)
            
        Raises:
            ValueError: If file validation fails
            Exception: If processing fails
        """
        logger.info(f"Processing file: {filename} ({file_type})")
        
        # Step 1: Validate
        is_valid, error_msg = self.validator.validate_file(filename, len(file_content))
        if not is_valid:
            logger.error(f"❌ File validation failed: {error_msg}")
            raise ValueError(error_msg)
        
        # Step 2: Save to disk
        file_path, file_id = self.storage.save_uploaded_file(
            file_content, filename, collection_id
        )
        
        try:
            # Step 3: Extract text
            full_text, page_count, page_details = self.extractor.extract_text(file_path, file_type)
            total_chars = len(full_text)
            
            # Step 4: Chunk text
            chunks = self.chunker.chunk_text(full_text)
            
            # Set file_id for all chunks
            for chunk in chunks:
                chunk.file_id = file_id
            
            # Try to detect page numbers
            if page_details and file_type in ["pdf", "pptx"]:
                page_markers = [f"--- Page {i + 1} ---" for i in range(len(page_details))]
                chunks = self.chunker.detect_page_numbers(chunks, page_markers)
            
            # Step 5: Create metadata
            metadata = StoredFileMetadata(
                file_id=file_id,
                collection_id=collection_id,
                filename=filename,
                file_type=file_type,
                file_size_bytes=len(file_content),
                file_path=file_path,
                pages=page_count if file_type in ["pdf", "pptx"] else 0,
                chunks=len(chunks),
                total_chars=total_chars,
                embedding_model=settings.embedding_model,
                uploaded_at=datetime.utcnow(),
                processed_at=datetime.utcnow(),
                status="processed"
            )
            
            logger.info(f"✅ File processing complete: {len(chunks)} chunks created")
            return chunks, metadata
            
        except Exception as e:
            # Clean up on failure
            self.storage.delete_file(file_path)
            logger.error(f"❌ File processing failed: {str(e)}")
            raise Exception(f"Failed to process file: {str(e)}")
        # At the END of file_processor.py (after all classes)

__all__ = ['FileProcessor', 'FileValidator', 'TextExtractor', 'TextChunker', 'FileStorage']